#!/usr/bin/env python3
"""Append a closing thank-you slide with team names and emails to the existing deck.

Only appends; existing slides are untouched.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
DECK = HERE / "Hunting-Leopards-with-Half-a-Brain.pptx"

PRIMARY = RGBColor(0xB4, 0x79, 0x54)
ACCENT = RGBColor(0x95, 0x5E, 0x39)
BACKGROUND = RGBColor(0xF7, 0xF1, 0xE8)
HEADING = RGBColor(0x54, 0x3E, 0x2D)
TEXT = RGBColor(0x3B, 0x2E, 0x23)
SUBTEXT = RGBColor(0x77, 0x63, 0x53)
FONT = "Trebuchet MS"

TEAM = [
    ("Omar Alharbi", "omar.alharbi@kaust.edu.sa"),
    ("Abdullah Alhindi", "abdullah.alhindi@kaust.edu.sa"),
    ("Abdullah Alghanim", "abdullah.alghanim@kaust.edu.sa"),
    ("Rasheed Hamidaddin", "rasheed.hamidaddin@kaust.edu.sa"),
    ("Basil Alshareef", "basilahmed.alshareef@gmail.com"),
    ("Basil Alshehri", "basil.alshehri55@gmail.com"),
]


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    return frame


def line(frame, text, size, color, bold=False, first=False, space_before=None):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    if space_before is not None:
        para.space_before = space_before
    run = para.add_run()
    run.text = text
    run.font.size = size
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = color
    return para


def main():
    prs = Presentation(DECK)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    rect(slide, 0, 0, prs.slide_width, prs.slide_height, BACKGROUND)
    rect(slide, 0, 0, Emu(310896), prs.slide_height, PRIMARY)
    rect(slide, 0, Emu(6547104), prs.slide_width, Emu(310896), PRIMARY)

    kicker = textbox(slide, Inches(1.3), Inches(1.05), Inches(10.4), Inches(0.4))
    line(kicker, "THE TEAM", Pt(13), ACCENT, bold=True, first=True)

    title = textbox(slide, Inches(1.3), Inches(1.5), Inches(10.4), Inches(0.9))
    line(title, "Thank You", Pt(44), HEADING, bold=True, first=True)

    rect(slide, Inches(1.3), Inches(2.55), Inches(1.6), Emu(50292), PRIMARY)

    col_x = (Inches(1.3), Inches(6.85))
    row_y = Inches(3.1)
    row_step = Inches(1.05)
    for index, (name, email) in enumerate(TEAM):
        left = col_x[index % 2]
        top = row_y + row_step * (index // 2)
        frame = textbox(slide, left, top, Inches(5.0), Inches(0.85))
        line(frame, name, Pt(19), TEXT, bold=True, first=True)
        line(frame, email, Pt(14), SUBTEXT)

    footer = textbox(slide, Inches(1.3), Inches(6.35), Inches(10.4), Inches(0.4))
    line(footer, "KAUST Academy", Pt(14), ACCENT, bold=True, first=True)

    prs.save(DECK)
    print(f"appended slide {len(prs.slides.__iter__.__self__._sldIdLst)} -> {DECK}")


if __name__ == "__main__":
    main()
