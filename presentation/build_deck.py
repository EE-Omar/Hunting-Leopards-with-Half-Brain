#!/usr/bin/env python3
"""Build the 5-minute presentation deck for Hunting Leopards with Half a Brain.

Extracts the figures it needs from the report PDF on first run, then writes the deck.
"""

import re
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).parent
REPO = HERE.parent
ASSETS = HERE / "assets"
OUT = HERE / "Hunting-Leopards-with-Half-a-Brain.pptx"

REPORT = Path("/home/abd0_linux/Sync/University/KAUST SUMMER/Final_Report_Leopard.pdf")

# Physical PDF pages, which run one ahead of the printed page numbers because of the title page.
ARCH_PAGE = 9   # printed page 8, Figure 1, vector so it has to be rendered and cropped
IR_PAGE = 16    # printed page 15, Figure 4, an embedded raster we can pull out directly

# Figure 1's bounding box on the rendered page, as fractions of page width and height.
ARCH_CROP = (0.09, 0.560, 0.91, 0.838)

PRIMARY = RGBColor(0xB4, 0x79, 0x54)      # brand copper, used for fills and markers
ACCENT = RGBColor(0x95, 0x5E, 0x39)       # deeper copper, legible as small text on cream
BACKGROUND = RGBColor(0xF7, 0xF1, 0xE8)   # warm cream
HEADING = RGBColor(0x54, 0x3E, 0x2D)      # deep brown
TEXT = RGBColor(0x3B, 0x2E, 0x23)         # body copy
SUBTEXT = RGBColor(0x77, 0x63, 0x53)      # muted brown for sub-bullets
LIGHT = RGBColor(0xFC, 0xFC, 0xFC)        # white, for figure mattes

FONT_HEAD = "Trebuchet MS"
FONT_BODY = "Trebuchet MS"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
BODY_TOP = Inches(2.05)


def trim_white(img, tolerance=12):
    background = Image.new("RGB", img.size, (255, 255, 255))
    diff = ImageChops.difference(img.convert("RGB"), background)
    box = ImageChops.add(diff, diff, 2.0, -tolerance).getbbox()
    return img.crop(box) if box else img


def extract_architecture(target):
    subprocess.run(
        ["pdftoppm", "-r", "300", "-f", str(ARCH_PAGE), "-l", str(ARCH_PAGE), "-png",
         str(REPORT), str(ASSETS / "_page")],
        check=True,
    )
    rendered = next(ASSETS.glob("_page*.png"))
    with Image.open(rendered) as page:
        width, height = page.size
        left, top, right, bottom = ARCH_CROP
        box = (int(left * width), int(top * height), int(right * width), int(bottom * height))
        trim_white(page.crop(box)).save(target)
    rendered.unlink()


def extract_ir_grid(target):
    subprocess.run(
        ["pdfimages", "-png", "-f", str(IR_PAGE), "-l", str(IR_PAGE),
         str(REPORT), str(ASSETS / "_ir")],
        check=True,
    )
    produced = sorted(ASSETS.glob("_ir*.png"))

    def area(path):
        with Image.open(path) as img:
            return img.size[0] * img.size[1]

    # The photo grid is by far the largest thing on the page; the rest is the header logo
    # and the alpha masks that pdfimages splits out alongside each image.
    biggest = max(produced, key=area)
    with Image.open(biggest) as img:
        img.convert("RGB").save(target)
    for path in produced:
        path.unlink()


def ensure_assets():
    ASSETS.mkdir(exist_ok=True)

    architecture = ASSETS / "architecture.png"
    if not architecture.exists():
        extract_architecture(architecture)

    ir_grid = ASSETS / "ir_detections.png"
    if not ir_grid.exists():
        extract_ir_grid(ir_grid)

    panels = []
    for name, source in (
        ("transmission_time.png", "results/charts/v3_panels/b_transmission_time.png"),
        ("accuracy.png", "results/charts/v3_panels/e_accuracy.png"),
    ):
        target = ASSETS / name
        if not target.exists():
            shutil.copy(REPO / source, target)
        panels.append(target)

    ir_strip = ASSETS / "ir_strip.png"
    if not ir_strip.exists():
        with Image.open(ir_grid) as img:
            width, height = img.size
            img.crop((0, 0, width, height // 4)).save(ir_strip)

    for path in (architecture, ir_grid, ir_strip, *panels):
        with Image.open(path) as img:
            print(f"  {path.name}: {img.size[0]}x{img.size[1]}")


def fill_background(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    return frame


def add_title(slide, text, eyebrow=None):
    if eyebrow:
        frame = add_textbox(slide, MARGIN, Inches(0.52), SLIDE_W - 2 * MARGIN, Inches(0.4))
        run = frame.paragraphs[0].add_run()
        run.text = eyebrow.upper()
        run.font.name = FONT_HEAD
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = ACCENT
        top = Inches(0.92)
    else:
        top = Inches(0.7)

    frame = add_textbox(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, Inches(0.9))
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = FONT_HEAD
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = HEADING

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN, Inches(1.78), Inches(1.6), Inches(0.055)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_bullets(slide, bullets, left, top, width, height, size=19):
    frame = add_textbox(slide, left, top, width, height)
    for index, (level, text) in enumerate(bullets):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.space_after = Pt(11 if level == 0 else 5)
        para.line_spacing = 1.15
        if level == 0:
            marker = para.add_run()
            marker.text = "▪  "
            marker.font.name = FONT_BODY
            marker.font.size = Pt(size)
            marker.font.color.rgb = PRIMARY
            body = para.add_run()
            body.text = text
            body.font.size = Pt(size)
            body.font.color.rgb = TEXT
        else:
            para.level = 1
            body = para.add_run()
            body.text = "      " + text
            body.font.size = Pt(size - 3)
            body.font.color.rgb = SUBTEXT
        body.font.name = FONT_BODY
    return frame


def add_stat(slide, left, top, width, number, label, detail):
    frame = add_textbox(slide, left, top, width, Inches(1.9))
    for index, (text, size, color, bold) in enumerate([
        (number, 80, ACCENT, True),
        (label, 18, TEXT, True),
        (detail, 15, SUBTEXT, False),
    ]):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = PP_ALIGN.CENTER
        para.space_after = Pt(4)
        run = para.add_run()
        run.text = text
        run.font.name = FONT_HEAD if index == 0 else FONT_BODY
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_image(slide, path, left, top, width, height, matte=True):
    """Fit an image inside the given box, preserving aspect ratio and centring it."""
    with Image.open(path) as img:
        img_w, img_h = img.size
    scale = min(width / img_w, height / img_h)
    draw_w, draw_h = int(img_w * scale), int(img_h * scale)
    x = int(left + (width - draw_w) / 2)
    y = int(top + (height - draw_h) / 2)

    if matte:
        pad = Inches(0.09)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Emu(x - pad), Emu(y - pad), Emu(draw_w + 2 * pad), Emu(draw_h + 2 * pad),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.fill.background()
        card.shadow.inherit = False

    slide.shapes.add_picture(str(path), Emu(x), Emu(y), Emu(draw_w), Emu(draw_h))


def new_slide(prs, background=BACKGROUND):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, background)
    return slide


def load_script():
    """Parse script.md into {slide number: spoken text} so the notes cannot drift from it."""
    text = (HERE / "script.md").read_text()
    parts = re.split(r"^## Slide (\d+):.*$", text, flags=re.M)[1:]
    sections = {}
    for number, body in zip(parts[::2], parts[1::2]):
        body = "\n".join(line for line in body.splitlines() if line.strip() != "---")
        sections[int(number)] = body.strip()
    return sections


SCRIPT = load_script()


def set_notes(slide, number):
    slide.notes_slide.notes_text_frame.text = SCRIPT[number]


def build():
    print("assets:")
    ensure_assets()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Title
    slide = new_slide(prs)
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.34), SLIDE_H
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()
    accent.shadow.inherit = False

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, SLIDE_H - Inches(0.34), SLIDE_W, Inches(0.34)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = PRIMARY
    band.line.fill.background()
    band.shadow.inherit = False

    frame = add_textbox(slide, Inches(1.3), Inches(1.95), Inches(11.0), Inches(1.9))
    for index, text in enumerate(["Hunting Leopards", "with Half a Brain"]):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = FONT_HEAD
        run.font.size = Pt(52)
        run.font.bold = True
        run.font.color.rgb = HEADING

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(4.02), Inches(2.4), Inches(0.06)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = PRIMARY
    rule.line.fill.background()
    rule.shadow.inherit = False

    frame = add_textbox(slide, Inches(1.3), Inches(4.25), Inches(10.5), Inches(0.5))
    run = frame.paragraphs[0].add_run()
    run.text = "Split inference for Arabian leopard detection on low power, off-grid hardware"
    run.font.name = FONT_BODY
    run.font.size = Pt(21)
    run.font.color.rgb = ACCENT

    frame = add_textbox(slide, Inches(1.3), Inches(5.35), Inches(11.0), Inches(1.2))
    for index, text in enumerate([
        "Omar Alharbi, Abdullah Alhindi, Abdullah Alghanim, Rasheed Hamidaddin,",
        "Basil Alshareef, Basil Alshehri, Mojtaba Alshams",
    ]):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = para.add_run()
        run.text = text
        run.font.name = FONT_BODY
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT
    para = frame.add_paragraph()
    para.space_before = Pt(10)
    run = para.add_run()
    run.text = "KAUST Academy"
    run.font.name = FONT_BODY
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    set_notes(slide, 1)

    # 2. Abstract
    slide = new_slide(prs)
    add_title(slide, "What we built", eyebrow="Abstract")
    add_bullets(slide, [
        (0, "One YOLO26 Large model, cut in two: the first three layers on a small camera node, "
            "everything after them on one shared server nearby."),
        (0, "A learned compression step sits right at the cut and shrinks what crosses the link "
            "by 64 times."),
        (0, "A large model runs on hardware that could never hold it, for almost no accuracy cost."),
    ], MARGIN, BODY_TOP, SLIDE_W - 2 * MARGIN, Inches(4.4), size=22)
    set_notes(slide, 2)

    # 3. Introduction
    slide = new_slide(prs)
    add_title(slide, "Why the Arabian leopard", eyebrow="Introduction")
    add_bullets(slide, [
        (0, "Roughly 200 left in the wild. Critically endangered on the IUCN Red List."),
        (0, "Saudi Arabia is investing heavily in its recovery:"),
        (1, "The Arabian Leopard Fund, run by the Royal Commission for AlUla"),
        (1, "Breeding at the Prince Saud Al-Faisal Wildlife Research Center in Taif, six cubs born in 2025"),
        (1, "Ibex, gazelle and oryx reintroduced to rebuild the food chain"),
        (1, "Partners include the IUCN, Panthera, Catmosphere and the Smithsonian"),
        (0, "All of it depends on one input: where the animals actually are, right now."),
        (0, "And that has to be collected in mountains with no power grid and no network."),
    ], MARGIN, BODY_TOP, Inches(11.6), Inches(4.6), size=20)
    set_notes(slide, 3)

    # 4. System Overview
    slide = new_slide(prs)
    add_title(slide, "A network of camera nodes, one shared server", eyebrow="System Overview")
    add_image(slide, ASSETS / "architecture.png",
              MARGIN, Inches(2.0), SLIDE_W - 2 * MARGIN, Inches(3.55))
    frame = add_textbox(slide, MARGIN, Inches(5.85), SLIDE_W - 2 * MARGIN, Inches(1.0))
    frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = frame.paragraphs[0].add_run()
    run.text = ("The heavy part of the model exists once and is shared, so covering more ground "
                "costs one more camera, not one more server.")
    run.font.name = FONT_BODY
    run.font.size = Pt(18)
    run.font.color.rgb = TEXT
    set_notes(slide, 4)

    # 5. Methodology: the split
    slide = new_slide(prs)
    add_title(slide, "Where we cut the model", eyebrow="Methodology")
    add_bullets(slide, [
        (0, "Data from LILA BC, seven open Roboflow datasets, and our own background images."),
        (0, "Eight classes. The leopard is the only one we act on."),
        (1, "Cheetah, hyena, Nubian ibex, camel, cat, dog, person, all labelled as hard negatives"),
        (0, "Cut point chosen empirically: 18 admissible candidates, all profiled on the real devices."),
        (0, "Layer 3 wins. It leaves 3.4% of the parameters on the camera, plus a small encoder."),
    ], MARGIN, BODY_TOP, Inches(11.6), Inches(4.5), size=20)
    set_notes(slide, 5)

    # 6. Methodology: the bottleneck
    slide = new_slide(prs)
    add_title(slide, "The learned bottleneck", eyebrow="Methodology")

    add_stat(slide, MARGIN, Inches(2.05), Inches(2.85),
             "64x", "smaller payload", "1600 KB to 25 KB per frame")
    add_stat(slide, Inches(3.9), Inches(2.05), Inches(2.85),
             "7x", "faster to send", "522 ms to 73 ms per frame")

    add_bullets(slide, [
        (0, "A 4,112 parameter encoder narrows the feature map from 256 channels to 16, sent as uint8."),
        (0, "A larger decoder rebuilds it on the server, where the compute is affordable."),
        (0, "Everything downstream of the split is left completely unchanged."),
    ], MARGIN, Inches(4.25), Inches(5.9), Inches(2.6), size=17)

    add_image(slide, ASSETS / "transmission_time.png",
              Inches(7.0), Inches(2.05), Inches(5.5), Inches(4.7))
    set_notes(slide, 6)

    # 7. Results
    slide = new_slide(prs)
    add_title(slide, "Accuracy holds, throughput improves", eyebrow="Results and Discussion")
    add_bullets(slide, [
        (0, "Leopard AP@0.5 of 0.988. Compression costs it 0.011 AP50-95."),
        (0, "Across all eight classes the drop is 0.0197 mAP50-95."),
        (0, "Measured on the deployed ONNX artefacts over a real link, not simulated."),
        (0, "Overlapping node and server, throughput rises 0.317 to 0.499 FPS, a 58% gain."),
    ], MARGIN, BODY_TOP, Inches(5.8), Inches(3.0), size=17)
    add_image(slide, ASSETS / "accuracy.png",
              Inches(6.9), Inches(1.95), Inches(5.6), Inches(3.35))
    add_image(slide, ASSETS / "ir_strip.png",
              MARGIN, Inches(5.65), Inches(11.6), Inches(1.6))
    set_notes(slide, 7)

    # 8. Conclusion
    slide = new_slide(prs)
    add_title(slide, "Conclusion", eyebrow="Conclusion")
    add_bullets(slide, [
        (0, "A model no single field device can run in real time now runs across a whole reserve."),
        (0, "The wireless link is no longer the limit on how many cameras a reserve can carry."),
        (0, "Next: skipping empty frames on the camera, detection-aware bottleneck training, a live field trial."),
    ], MARGIN, BODY_TOP, Inches(11.6), Inches(3.0), size=22)

    frame = add_textbox(slide, MARGIN, Inches(5.5), Inches(11.6), Inches(0.9))
    run = frame.paragraphs[0].add_run()
    run.text = "github.com/EE-Omar/Hunting-Leopards-with-Half-Brain"
    run.font.name = FONT_BODY
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    set_notes(slide, 8)

    check_no_em_dashes(prs)
    prs.save(OUT)
    print(f"wrote {OUT.name} ({len(prs.slides)} slides)")


def check_no_em_dashes(prs):
    for number, slide in enumerate(prs.slides, start=1):
        texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
        if slide.has_notes_slide:
            texts.append(slide.notes_slide.notes_text_frame.text)
        for text in texts:
            if "—" in text or "–" in text:
                raise SystemExit(f"slide {number} contains a dash character: {text!r}")


if __name__ == "__main__":
    build()
